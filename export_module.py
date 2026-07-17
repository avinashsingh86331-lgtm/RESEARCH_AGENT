"""
export_module.py
-----------------
Job: Turn the AI's output into two real files:
   1. A PowerPoint (.pptx) - using python-pptx
   2. A PDF report (.pdf)  - using fpdf2

Both libraries are pure Python, need no extra software (no MS Office,
no LibreOffice) and no database.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF
from fpdf.enums import XPos, YPos


def _make_wrappable(text, max_word_len=40):
    """
    fpdf2's multi_cell() wraps at spaces. A long unbroken string (like a
    URL with no spaces) that's wider than the page has nowhere to wrap,
    which raises FPDFException: 'Not enough horizontal space to render
    a single character'. This inserts a zero-width break opportunity
    (a soft space every max_word_len characters) inside any "word"
    longer than that, so multi_cell always has somewhere to wrap.
    """
    words = text.split(" ")
    fixed_words = []
    for word in words:
        if len(word) > max_word_len:
            # Break the long word into chunks, joined with a real space
            # so multi_cell can wrap between them.
            chunks = [word[i:i + max_word_len] for i in range(0, len(word), max_word_len)]
            fixed_words.append(" ".join(chunks))
        else:
            fixed_words.append(word)
    return " ".join(fixed_words)


def create_pptx(topic, slide_outline, output_path="research_output.pptx"):
    """Build a simple, clean slide deck from the AI's slide_outline JSON."""
    prs = Presentation()

    # --- Title slide ---
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = slide_outline.get("title_slide", topic)
    slide.placeholders[1].text = slide_outline.get("subtitle", "")

    # --- Content slides ---
    bullet_layout = prs.slide_layouts[1]
    for item in slide_outline.get("slides", []):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = item.get("heading", "")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = item.get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)

    prs.save(output_path)
    return output_path


def create_pdf(topic, notes_text, source_list, output_path="research_output.pdf"):
    """Build a PDF report containing the notes and a numbered reference list."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Title
    pdf.set_font("Helvetica", "B", 18)
    pdf.multi_cell(0, 10, f"Research Report: {topic}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(4)

    # Notes body
    pdf.set_font("Helvetica", "", 12)
    for line in notes_text.split("\n"):
        clean_line = line.encode("latin-1", "replace").decode("latin-1")
        clean_line = _make_wrappable(clean_line)
        if clean_line.strip().startswith("#"):
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(0, 8, clean_line.replace("#", "").strip(), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Helvetica", "", 12)
        else:
            pdf.multi_cell(0, 7, clean_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)

    # References section
    pdf.ln(6)
    pdf.set_font("Helvetica", "B", 14)
    pdf.multi_cell(0, 8, "References", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.set_font("Helvetica", "", 11)
    for s in source_list:
        ref_line = f"[{s['num']}] {s['title']} - {s['url']}"
        ref_line = ref_line.encode("latin-1", "replace").decode("latin-1")
        ref_line = _make_wrappable(ref_line)
        pdf.multi_cell(0, 7, ref_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.output(output_path)
    return output_path